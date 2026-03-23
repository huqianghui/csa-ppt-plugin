#!/usr/bin/env python3
"""Core local PPT generation, template, and edit helpers."""

from __future__ import annotations

import json
import os
import re
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import DEFAULT_MAX_SLIDES, default_slide_theme


def _hex_to_rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    if len(value) != 6:
        value = "0F6CBD"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _remove_slide(prs: Presentation, index: int) -> None:
    slide_ids = prs.slides._sldIdLst
    slide_elem = slide_ids[index]
    rel_id = slide_elem.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    slide_ids.remove(slide_elem)
    prs.part.rels.pop(rel_id)


def _clear_text_frame(text_frame) -> None:
    text_frame.clear()


def _set_run_style(run, font_size: int, color: RGBColor, bold: bool = False) -> None:
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_shape_text(shape, title: str, font_size: int, color: RGBColor, bold: bool = False) -> None:
    text_frame = shape.text_frame
    _clear_text_frame(text_frame)
    paragraph = text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.alignment = PP_ALIGN.LEFT
    if paragraph.runs:
        _set_run_style(paragraph.runs[0], font_size, color, bold=bold)


def _find_placeholder(slide, placeholder_type: Optional[int] = None):
    for placeholder in slide.placeholders:
        try:
            current_type = placeholder.placeholder_format.type
        except Exception:
            continue
        if placeholder_type is None or current_type == placeholder_type:
            return placeholder
    return None


def _find_body_placeholder(slide):
    for placeholder_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE):
        placeholder = _find_placeholder(slide, placeholder_type)
        if placeholder is not None:
            return placeholder
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape != slide.shapes.title:
            return shape
    return None


def _pick_layout(prs: Presentation, role: str):
    desired_names = {
        "title": ["title", "cover"],
        "content": ["title and content", "content", "title only"],
        "section": ["section", "title only"],
        "blank": ["blank"],
    }
    candidates = desired_names.get(role, [])
    for wanted in candidates:
        for layout in prs.slide_layouts:
            if wanted in layout.name.lower():
                return layout
    if role == "blank":
        return prs.slide_layouts[-1]
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]


def _add_textbox(slide, left: float, top: float, width: float, height: float):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _write_bullets(shape, bullets: Iterable[str], color: RGBColor, font_size: int = 20) -> None:
    text_frame = shape.text_frame
    _clear_text_frame(text_frame)
    bullets = list(bullets)
    if not bullets:
        bullets = ["Fill in the main point here."]
    first = text_frame.paragraphs[0]
    first.text = bullets[0]
    first.level = 0
    if first.runs:
        _set_run_style(first.runs[0], font_size, color)
    for item in bullets[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        if paragraph.runs:
            _set_run_style(paragraph.runs[0], font_size, color)


def _apply_background(slide, theme: Dict[str, str]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(theme.get("background", "#FFFFFF"))


def _render_title_slide(prs: Presentation, slide_spec: Dict[str, Any], theme: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs, "title"))
    _apply_background(slide, theme)
    title_color = _hex_to_rgb(theme.get("accent", "#0F6CBD"))
    text_color = _hex_to_rgb(theme.get("text", "#1F1F1F"))

    title_shape = slide.shapes.title
    if title_shape is None:
        title_shape = _add_textbox(slide, 0.8, 1.2, 11.0, 1.4)
    _set_shape_text(title_shape, slide_spec.get("title", "Untitled deck"), 30, title_color, bold=True)

    subtitle = slide_spec.get("subtitle") or slide_spec.get("summary") or slide_spec.get("query", "")
    body_shape = _find_body_placeholder(slide)
    if body_shape is None:
        body_shape = _add_textbox(slide, 0.9, 3.0, 10.6, 1.8)
    _set_shape_text(body_shape, subtitle, 18, text_color)


def _render_bullets_slide(prs: Presentation, slide_spec: Dict[str, Any], theme: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs, "content"))
    _apply_background(slide, theme)
    title_color = _hex_to_rgb(theme.get("accent", "#0F6CBD"))
    text_color = _hex_to_rgb(theme.get("text", "#1F1F1F"))

    title_shape = slide.shapes.title or _add_textbox(slide, 0.7, 0.5, 11.2, 0.8)
    _set_shape_text(title_shape, slide_spec.get("title", "Slide"), 24, title_color, bold=True)

    body_shape = _find_body_placeholder(slide)
    if body_shape is None:
        body_shape = _add_textbox(slide, 0.9, 1.5, 10.6, 4.8)
    bullets = slide_spec.get("bullets") or slide_spec.get("points") or []
    _write_bullets(body_shape, bullets, text_color)


def _render_section_slide(prs: Presentation, slide_spec: Dict[str, Any], theme: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs, "section"))
    _apply_background(slide, theme)
    title_color = _hex_to_rgb(theme.get("accent", "#0F6CBD"))
    text_color = _hex_to_rgb(theme.get("muted", "#5B5B5B"))
    title_shape = slide.shapes.title or _add_textbox(slide, 0.8, 1.5, 11.0, 1.2)
    _set_shape_text(title_shape, slide_spec.get("title", "Section"), 28, title_color, bold=True)
    subtitle_shape = _find_body_placeholder(slide)
    if subtitle_shape is None:
        subtitle_shape = _add_textbox(slide, 0.9, 3.0, 10.4, 1.4)
    _set_shape_text(subtitle_shape, slide_spec.get("subtitle", ""), 18, text_color)


def _render_quote_slide(prs: Presentation, slide_spec: Dict[str, Any], theme: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs, "blank"))
    _apply_background(slide, theme)
    title_color = _hex_to_rgb(theme.get("accent", "#0F6CBD"))
    text_color = _hex_to_rgb(theme.get("text", "#1F1F1F"))

    title_box = _add_textbox(slide, 0.8, 0.6, 11.0, 0.8)
    _set_shape_text(title_box, slide_spec.get("title", "Quote"), 24, title_color, bold=True)

    quote_box = _add_textbox(slide, 1.0, 1.8, 10.2, 2.6)
    _set_shape_text(quote_box, slide_spec.get("quote", ""), 24, text_color, bold=False)

    attribution_box = _add_textbox(slide, 1.0, 4.8, 10.0, 0.7)
    _set_shape_text(attribution_box, slide_spec.get("attribution", ""), 16, text_color)


def _render_two_column_slide(prs: Presentation, slide_spec: Dict[str, Any], theme: Dict[str, str]) -> None:
    slide = prs.slides.add_slide(_pick_layout(prs, "blank"))
    _apply_background(slide, theme)
    title_color = _hex_to_rgb(theme.get("accent", "#0F6CBD"))
    text_color = _hex_to_rgb(theme.get("text", "#1F1F1F"))

    title_box = _add_textbox(slide, 0.7, 0.5, 11.2, 0.8)
    _set_shape_text(title_box, slide_spec.get("title", "Comparison"), 24, title_color, bold=True)

    left_box = _add_textbox(slide, 0.8, 1.6, 5.1, 4.8)
    left_title = slide_spec.get("left_title") or "Left"
    left_bullets = [left_title] + list(slide_spec.get("left_bullets") or [])
    _write_bullets(left_box, left_bullets, text_color)

    right_box = _add_textbox(slide, 6.2, 1.6, 5.1, 4.8)
    right_title = slide_spec.get("right_title") or "Right"
    right_bullets = [right_title] + list(slide_spec.get("right_bullets") or [])
    _write_bullets(right_box, right_bullets, text_color)


def normalize_outline(outline: Dict[str, Any], query: str = "", language: str = "Chinese") -> Dict[str, Any]:
    normalized = dict(outline or {})
    normalized.setdefault("title", outline.get("title") if outline else query or "Presentation")
    normalized.setdefault("language", outline.get("language") if outline else language)
    normalized.setdefault("theme", dict(default_slide_theme(), **(outline.get("theme") if outline else {})))
    slides = list(normalized.get("slides") or [])
    if not slides:
        slides = build_basic_outline(query or normalized["title"], "", language).get("slides", [])
    normalized["slides"] = slides
    return normalized


def render_presentation(outline: Dict[str, Any], output_path: str, template_file: str = "") -> Dict[str, Any]:
    prs = Presentation(template_file) if template_file else Presentation()
    while len(prs.slides) > 0:
        _remove_slide(prs, 0)

    normalized = normalize_outline(outline)
    theme = normalized.get("theme") or default_slide_theme()
    rendered_titles: List[str] = []
    for slide_spec in normalized.get("slides", []):
        slide_type = (slide_spec.get("type") or "bullets").lower()
        rendered_titles.append(slide_spec.get("title", slide_type.title()))
        if slide_type == "title":
            _render_title_slide(prs, slide_spec, theme)
        elif slide_type == "section":
            _render_section_slide(prs, slide_spec, theme)
        elif slide_type == "quote":
            _render_quote_slide(prs, slide_spec, theme)
        elif slide_type == "two-column":
            _render_two_column_slide(prs, slide_spec, theme)
        else:
            _render_bullets_slide(prs, slide_spec, theme)

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {
        "output_file": str(output.resolve()),
        "slide_count": len(prs.slides),
        "slide_titles": rendered_titles,
    }


def _clean_line(line: str) -> str:
    return re.sub(r"\s+", " ", line.strip())


def _extract_bullet_candidates(reference_text: str) -> List[str]:
    bullets: List[str] = []
    for raw_line in reference_text.splitlines():
        line = _clean_line(raw_line)
        if not line:
            continue
        if line.startswith(("- ", "* ")):
            bullets.append(line[2:].strip())
        elif re.match(r"^\d+[\.)]\s+", line):
            bullets.append(re.sub(r"^\d+[\.)]\s+", "", line))
    return bullets


def _extract_sections(reference_text: str) -> List[Tuple[str, List[str]]]:
    sections: List[Tuple[str, List[str]]] = []
    current_title = "Overview"
    current_items: List[str] = []
    for raw_line in reference_text.splitlines():
        line = _clean_line(raw_line)
        if not line:
            continue
        if line.startswith("#"):
            if current_items:
                sections.append((current_title, current_items))
            current_title = line.lstrip("# ")
            current_items = []
            continue
        if line.startswith(("- ", "* ")):
            current_items.append(line[2:].strip())
        elif re.match(r"^\d+[\.)]\s+", line):
            current_items.append(re.sub(r"^\d+[\.)]\s+", "", line))
        elif len(current_items) < 4:
            current_items.append(line)
    if current_items:
        sections.append((current_title, current_items))
    return sections


def build_basic_outline(query: str, reference_text: str, language: str, max_slides: int = DEFAULT_MAX_SLIDES) -> Dict[str, Any]:
    title = query.strip() or "Presentation"
    sections = _extract_sections(reference_text)
    bullets = _extract_bullet_candidates(reference_text)
    if not sections and bullets:
        sections = [("Key Points", bullets[:8])]

    slides: List[Dict[str, Any]] = [
        {
            "type": "title",
            "title": title,
            "subtitle": f"Generated locally for {language} delivery",
        }
    ]

    if bullets:
        slides.append(
            {
                "type": "bullets",
                "title": "Executive Summary",
                "bullets": bullets[:4],
            }
        )
    elif reference_text.strip():
        intro_lines = [_clean_line(line) for line in reference_text.splitlines() if _clean_line(line)]
        slides.append(
            {
                "type": "bullets",
                "title": "Overview",
                "bullets": intro_lines[:4],
            }
        )
    else:
        slides.append(
            {
                "type": "bullets",
                "title": "Overview",
                "bullets": [
                    "Objectives and audience",
                    "Current context",
                    "Recommendations",
                    "Next steps",
                ],
            }
        )

    for section_title, items in sections[: max(0, max_slides - 3)]:
        slides.append(
            {
                "type": "bullets",
                "title": section_title,
                "bullets": items[:5],
            }
        )

    slides = slides[: max(1, max_slides - 1)]
    slides.append(
        {
            "type": "bullets",
            "title": "Next Steps",
            "bullets": [
                "Confirm slide narrative and audience expectations",
                "Add any charts, screenshots, or speaker notes",
                "Review style consistency before delivery",
            ],
        }
    )

    return {
        "title": title,
        "language": language,
        "theme": default_slide_theme(),
        "slides": slides[:max_slides],
    }


def analyze_template(template_path: str) -> Dict[str, Any]:
    prs = Presentation(template_path)
    layouts = []
    for index, layout in enumerate(prs.slide_layouts):
        layouts.append(
            {
                "index": index,
                "name": layout.name,
                "placeholders": len(layout.placeholders),
            }
        )

    sample_titles = []
    for index, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        if not title:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = _clean_line(shape.text_frame.text)
                    if text:
                        title = text[:80]
                        break
        sample_titles.append({"slide": index, "title": title or "(no title)"})

    return {
        "template_file": str(Path(template_path).resolve()),
        "slide_count": len(prs.slides),
        "slide_width_inches": round(prs.slide_width.inches, 2),
        "slide_height_inches": round(prs.slide_height.inches, 2),
        "layouts": layouts,
        "sample_titles": sample_titles,
    }


def extract_presentation_text(pptx_path: str) -> Dict[str, Any]:
    prs = Presentation(pptx_path)
    slides = []
    all_text: List[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = _clean_line(shape.text_frame.text)
                if text:
                    texts.append(text)
                    all_text.append(text)
        slides.append({"slide": index, "texts": texts})
    return {
        "slide_count": len(slides),
        "slides": slides,
        "text": "\n".join(all_text),
    }


def load_json_file(file_path: str) -> Dict[str, Any]:
    with open(file_path, "r", encoding="utf-8") as handle:
        return json.load(handle)


def apply_edit_plan(input_file: str, output_file: str, edit_plan: Dict[str, Any]) -> Dict[str, Any]:
    prs = Presentation(input_file)
    summary: List[str] = []
    edits = list(edit_plan.get("edits") or [])
    for edit in edits:
        action = (edit.get("action") or "").lower()
        if action == "update_text":
            slide_number = int(edit["slide"])
            if slide_number < 1 or slide_number > len(prs.slides):
                raise ValueError(f"Slide {slide_number} is out of range for this presentation")
            slide = prs.slides[slide_number - 1]
            title = edit.get("title", "")
            bullets = edit.get("bullets") or []
            if slide.shapes.title is not None and title:
                _set_shape_text(slide.shapes.title, title, 24, _hex_to_rgb(default_slide_theme()["accent"]), bold=True)
            body_shape = _find_body_placeholder(slide)
            if body_shape is None:
                body_shape = _add_textbox(slide, 0.9, 1.6, 10.4, 4.8)
            if bullets:
                _write_bullets(body_shape, bullets, _hex_to_rgb(default_slide_theme()["text"]))
            summary.append(f"Updated slide {slide_number}")
        elif action == "append_slide":
            slide_spec = dict(edit)
            slide_spec.pop("action", None)
            slide_type = (slide_spec.get("type") or "bullets").lower()
            theme = default_slide_theme()
            if slide_type == "title":
                _render_title_slide(prs, slide_spec, theme)
            elif slide_type == "section":
                _render_section_slide(prs, slide_spec, theme)
            elif slide_type == "quote":
                _render_quote_slide(prs, slide_spec, theme)
            elif slide_type == "two-column":
                _render_two_column_slide(prs, slide_spec, theme)
            else:
                _render_bullets_slide(prs, slide_spec, theme)
            summary.append(f"Appended a new {slide_type} slide")
        elif action == "delete_slide":
            slide_number = int(edit["slide"])
            if slide_number < 1 or slide_number > len(prs.slides):
                raise ValueError(f"Slide {slide_number} is out of range for this presentation")
            _remove_slide(prs, slide_number - 1)
            summary.append(f"Deleted slide {slide_number}")
        elif action == "reorder_slides":
            order = [int(item) - 1 for item in edit.get("order") or []]
            total = len(prs.slides)
            if len(order) != total:
                raise ValueError("Reorder action must include every slide exactly once")
            if sorted(order) != list(range(total)):
                raise ValueError("Reorder action contains duplicate or missing slide numbers")
            slide_ids = prs.slides._sldIdLst
            existing = list(slide_ids)
            for elem in existing:
                slide_ids.remove(elem)
            for idx in order:
                slide_ids.append(existing[idx])
            summary.append("Reordered slides")
        else:
            raise ValueError(f"Unsupported edit action: {action}")

    output = Path(output_file)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {
        "output_file": str(output.resolve()),
        "slide_count": len(prs.slides),
        "summary": summary,
    }
